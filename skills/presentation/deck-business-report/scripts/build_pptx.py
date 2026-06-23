#!/usr/bin/env python3
"""Build a business/industrial report .pptx from an approved slide-spec.

Usage:
    pip install python-pptx
    python build_pptx.py slides.json [-o deck.pptx] [--template brand.pptx]

The input JSON mirrors the slide-spec the deck-business-report skill produces.
Each slide is rendered answer-first with a claim-title. Supported slide types:

    title    : {"type":"title","title":..., "subtitle":..., "author":..., "date":...}
    section  : {"type":"section","title":...}
    bullets  : {"type":"bullets","title":..., "bullets":[...], "notes":...}
    table    : {"type":"table","title":..., "columns":[...], "rows":[[...],...], "notes":...}
    chart    : {"type":"chart","title":..., "chart":"bar"|"line"|"pie",
                "categories":[...], "series":[{"name":...,"values":[...]}, ...], "notes":...}
    image    : {"type":"image","title":..., "image":"path.png" | null, "caption":..., "notes":...}

Top-level schema:
    {
      "theme": {"accent":"#1F6FEB", "text":"#1A1A1A", "muted":"#6E7781"},  # optional
      "footer": "Company — Confidential",   # optional, shown on content slides
      "slides": [ {slide}, ... ]
    }

Notes:
- bullets accept nested lists via {"text": "...", "level": 1} or plain strings.
- A 'chart' with no series renders a labeled placeholder box (paste a branded chart).
- An 'image' with null/missing path renders a TODO placeholder.
- Pass --template to inherit a corporate master (brand fonts/colors/logo).
"""
from __future__ import annotations

import argparse
import json
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
except ImportError:
    sys.exit("python-pptx is required: pip install python-pptx")


def hex_to_rgb(h: str) -> "RGBColor":
    return RGBColor.from_string(h.lstrip("#").upper())


DEFAULT_THEME = {"accent": "#1F6FEB", "text": "#1A1A1A", "muted": "#6E7781"}
# 16:9 canvas dimensions (used when not inheriting a template).
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
}


class DeckBuilder:
    def __init__(self, spec: dict, template: str | None):
        self.theme = {**DEFAULT_THEME, **spec.get("theme", {})}
        self.footer = spec.get("footer")
        self.prs = Presentation(template) if template else Presentation()
        if not template:
            self.prs.slide_width, self.prs.slide_height = SLIDE_W, SLIDE_H
        # Use the blank layout so we control placement explicitly.
        self.blank = self.prs.slide_layouts[6]

    # -- low-level helpers --------------------------------------------------
    def _slide(self):
        return self.prs.slides.add_slide(self.blank)

    def _claim_title(self, slide, text: str):
        box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35),
                                       self.prs.slide_width - Inches(1.2), Inches(1.1))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(26)
        run.font.bold = True
        run.font.color.rgb = hex_to_rgb(self.theme["accent"])
        return box

    def _notes(self, slide, notes: str | None):
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _footer(self, slide):
        if not self.footer:
            return
        box = slide.shapes.add_textbox(Inches(0.6), self.prs.slide_height - Inches(0.5),
                                       self.prs.slide_width - Inches(1.2), Inches(0.35))
        r = box.text_frame.paragraphs[0].add_run()
        r.text = self.footer
        r.font.size = Pt(9)
        r.font.color.rgb = hex_to_rgb(self.theme["muted"])

    def _placeholder(self, slide, label: str):
        from pptx.enum.shapes import MSO_SHAPE
        shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2.5), Inches(2.2),
                                     Inches(8.3), Inches(3.6))
        shp.fill.background()
        shp.line.color.rgb = hex_to_rgb(self.theme["muted"])
        tf = shp.text_frame
        tf.word_wrap = True
        r = tf.paragraphs[0].add_run()
        r.text = label
        r.font.size = Pt(14)
        r.font.color.rgb = hex_to_rgb(self.theme["muted"])

    # -- slide types --------------------------------------------------------
    def title(self, s: dict):
        slide = self._slide()
        box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4),
                                       self.prs.slide_width - Inches(1.6), Inches(2.5))
        tf = box.text_frame
        tf.word_wrap = True
        t = tf.paragraphs[0].add_run()
        t.text = s.get("title", "TODO: title")
        t.font.size = Pt(40)
        t.font.bold = True
        t.font.color.rgb = hex_to_rgb(self.theme["text"])
        for key, size in (("subtitle", 20), ("author", 16), ("date", 16)):
            if s.get(key):
                p = tf.add_paragraph()
                r = p.add_run()
                r.text = s[key]
                r.font.size = Pt(size)
                r.font.color.rgb = hex_to_rgb(self.theme["muted"])

    def section(self, s: dict):
        slide = self._slide()
        box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0),
                                       self.prs.slide_width - Inches(1.6), Inches(1.5))
        r = box.text_frame.paragraphs[0].add_run()
        r.text = s.get("title", "")
        r.font.size = Pt(30)
        r.font.bold = True
        r.font.color.rgb = hex_to_rgb(self.theme["accent"])

    def bullets(self, s: dict):
        slide = self._slide()
        self._claim_title(slide, s.get("title", ""))
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7),
                                       self.prs.slide_width - Inches(1.4), Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for b in s.get("bullets", []):
            text = b["text"] if isinstance(b, dict) else str(b)
            level = b.get("level", 0) if isinstance(b, dict) else 0
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.level = level
            r = p.add_run()
            r.text = ("• " if level == 0 else "– ") + text
            r.font.size = Pt(20 - 2 * min(level, 2))
            r.font.color.rgb = hex_to_rgb(self.theme["text"])
        self._footer(slide)
        self._notes(slide, s.get("notes"))

    def table(self, s: dict):
        slide = self._slide()
        self._claim_title(slide, s.get("title", ""))
        cols = s.get("columns", [])
        rows = s.get("rows", [])
        if not cols:
            self._placeholder(slide, "TODO: table data")
        else:
            n_rows, n_cols = len(rows) + 1, len(cols)
            tbl = slide.shapes.add_table(n_rows, n_cols, Inches(0.7), Inches(1.8),
                                         self.prs.slide_width - Inches(1.4),
                                         Inches(0.4) * n_rows).table
            for j, c in enumerate(cols):
                cell = tbl.cell(0, j)
                cell.text = str(c)
                cell.text_frame.paragraphs[0].runs[0].font.bold = True
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(13)
            for i, row in enumerate(rows, start=1):
                for j, val in enumerate(row):
                    cell = tbl.cell(i, j)
                    cell.text = str(val)
                    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
        self._footer(slide)
        self._notes(slide, s.get("notes"))

    def chart(self, s: dict):
        slide = self._slide()
        self._claim_title(slide, s.get("title", ""))
        series = s.get("series", [])
        if not series:
            self._placeholder(slide, "TODO: paste branded chart "
                                     f"({s.get('chart', 'bar')})")
        else:
            data = CategoryChartData()
            data.categories = s.get("categories", [])
            for ser in series:
                data.add_series(ser.get("name", ""), ser.get("values", []))
            ctype = CHART_TYPES.get(s.get("chart", "bar"), XL_CHART_TYPE.COLUMN_CLUSTERED)
            slide.shapes.add_chart(ctype, Inches(1.0), Inches(1.8),
                                   Inches(11.3), Inches(4.8), data)
        self._footer(slide)
        self._notes(slide, s.get("notes"))

    def image(self, s: dict):
        slide = self._slide()
        self._claim_title(slide, s.get("title", ""))
        path = s.get("image")
        if path:
            slide.shapes.add_picture(path, Inches(1.5), Inches(1.8), height=Inches(4.8))
        else:
            self._placeholder(slide, "TODO: image — " + s.get("caption", ""))
        self._footer(slide)
        self._notes(slide, s.get("notes"))

    # -- driver -------------------------------------------------------------
    def build(self, slides: list[dict]):
        dispatch = {
            "title": self.title, "section": self.section, "bullets": self.bullets,
            "table": self.table, "chart": self.chart, "image": self.image,
        }
        for i, s in enumerate(slides):
            kind = s.get("type", "bullets")
            handler = dispatch.get(kind)
            if handler is None:
                sys.exit(f"slide {i}: unknown type {kind!r}")
            handler(s)
        return self.prs


def main() -> None:
    ap = argparse.ArgumentParser(description="Build a business report .pptx from a slide-spec JSON.")
    ap.add_argument("spec", help="slides.json")
    ap.add_argument("-o", "--out", default="deck.pptx")
    ap.add_argument("--template", default=None, help="corporate .pptx template to inherit")
    args = ap.parse_args()

    with open(args.spec, encoding="utf-8") as f:
        spec = json.load(f)

    prs = DeckBuilder(spec, args.template).build(spec.get("slides", []))
    prs.save(args.out)
    print(f"wrote {args.out} ({len(spec.get('slides', []))} slides)")


if __name__ == "__main__":
    main()
